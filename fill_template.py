from genericpath import isdir
from os import listdir, path
import re
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Cm
from openpyxl import load_workbook

book = load_workbook(path.join("questions", "questions.xlsx"), read_only=True)

ROUNDS = [name.upper() for name in book.sheetnames if name != "Wipeout"]
assert len(ROUNDS) == 3, "Need to have 3 rounds in the excel workbook"
assert book.sheetnames[3] == "Wipeout", "4th sheet in excel is not called Wipeout"

prs = Presentation("template.pptx")
slide_layouts = prs.slide_master.slide_layouts


def create_rounds_preview():
    rounds_slide = prs.slides.add_slide(slide_layouts.get_by_name("ROUNDS"))
    Q1_idx, Q2_idx, Q3_idx, Q4_idx, Q5_idx, Q6_idx = 10, 11, 12, 13, 14, 15 # defined when creating the objects, can't change these easily
    rounds_slide.placeholders[Q1_idx].text = f"1. {ROUNDS[0]} - 10 QUESTIONS"
    rounds_slide.placeholders[Q2_idx].text = f"2. {ROUNDS[1]} - 10 QUESTIONS"
    rounds_slide.placeholders[Q3_idx].text = f"3. {ROUNDS[2]} - 10 QUESTIONS"
    rounds_slide.placeholders[Q4_idx].text = f"4. PICTURE ROUND - 9 QUESTIONS"
    rounds_slide.placeholders[Q5_idx].text = f"5. WIPEOUT ROUND - 10 QUESTIONS"
    rounds_slide.placeholders[Q6_idx].text = f"6. MUSIC ROUND - 10 QUESTIONS"


def create_worded_round(round_number):
    round_name = ROUNDS[round_number - 1]
    questions = [
        {"question": pair[0].value, "answer": pair[1].value}
        for pair in book.worksheets[round_number - 1]['A2':'B11']
        if pair[0].value is not None and pair[1].value is not None
    ]
    assert len(questions) == 10, f"Round {round_number} called {round_name} does not have 10 questions"
    round_intro_slide = prs.slides.add_slide(slide_layouts.get_by_name("ROUND_ANNOUNCEMENT"))
    round_intro_slide.placeholders[10].text = f"{round_name} ROUND"
    general_knowledge_slide = prs.slides.add_slide(slide_layouts.get_by_name("WORDED_QUESTIONS"))
    general_knowledge_slide.placeholders[22].text = f"{round_name} QUESTIONS"

    for idx, pair in enumerate(questions):
        general_knowledge_slide.placeholders[idx+33].text = f"{idx+1}. {pair['question']}"

    round_intro_slide = prs.slides.add_slide(slide_layouts.get_by_name("ROUND_ANNOUNCEMENT"))
    round_intro_slide.placeholders[10].text = f"{round_name} ANSWERS"

    for idx, pair in enumerate(questions):
        answer_slide = prs.slides.add_slide(slide_layouts.get_by_name("ROUND_ANSWER"))
        answer_slide.placeholders[22].text = f"{round_name} ANSWERS"
        answer_slide.placeholders[10].text = f"{idx+1}. {pair['question']}"
        answer_slide.placeholders[23].text = f"{idx+1}. {pair['answer']}"


def create_picture_round():
    picture_round_intro_slide = prs.slides.add_slide(slide_layouts.get_by_name("ROUND_ANNOUNCEMENT"))
    picture_round_intro_slide.placeholders[10].text = "PICTURE ROUND"
    picture_round_slide = prs.slides.add_slide(slide_layouts.get_by_name("PICTURE_QUESTIONS"))

    pictures = [pic for pic in listdir(path.join("questions", "pictures")) if re.match(".*\.(jpe?g|png)$", pic)]
    assert len(pictures) == 9, "Need to have 9 images in the questions/pictures folder"
    for idx, img_path in enumerate(pictures):
        picture_round_slide.placeholders[idx+10].insert_picture(path.join("questions", "pictures", img_path))
        picture_round_slide.placeholders[idx+19].text = f"{idx+1}."

    picture_round_break_slide = prs.slides.add_slide(slide_layouts.get_by_name("ROUND_ANNOUNCEMENT"))
    picture_round_break_slide.placeholders[10].text = "PICTURE ROUND ANSWERS"

    for idx, img_path in enumerate(pictures):
        picture_answer_slide = prs.slides.add_slide(slide_layouts.get_by_name("PICTURE_ANSWER"))
        picture_answer_slide.placeholders[22].text = f"PICTURE ROUND ANSWERS"
        picture_answer_slide.placeholders[11].text = f"{idx+1}."
        picture_answer_slide.placeholders[23].text = re.sub("\.jpeg|\.png", "", img_path)
        picture_answer_slide.placeholders[10].insert_picture(path.join("questions", "pictures", img_path))


def create_wipeout_round():
    WIPEOUT_ROUND_QS = [
        {
            "question": qa[0].value,
            "options": [cell.value for cell in qa[1:5]],
            "correct_option": qa[5].value
        }
        for qa in book.worksheets[3]["A2":"F11"]
        if len([cell for cell in qa if cell.value is not None]) == 6
    ]
    assert len(WIPEOUT_ROUND_QS) == 10, "There are not 10 complete wipeout questions"
    prs.slides.add_slide(slide_layouts.get_by_name("WIPEOUT_RULES"))
    for idx, question in enumerate(WIPEOUT_ROUND_QS):
        wipeout_question_slide = prs.slides.add_slide(slide_layouts.get_by_name("WIPEOUT_QUESTION"))
        wipeout_question_slide.placeholders[10].text = f"QUESTION {idx+1}."
        wipeout_question_slide.placeholders[11].text = question['question']
        wipeout_question_slide.placeholders[12].text = f"A. {question['options'][0]}"
        wipeout_question_slide.placeholders[13].text = f"B. {question['options'][1]}"
        wipeout_question_slide.placeholders[14].text = f"C. {question['options'][2]}"
        wipeout_question_slide.placeholders[15].text = f"D. {question['options'][3]}"
    prs.slides.add_slide(slide_layouts.get_by_name("WIPEOUT_ANSWER_BREAK"))
    for idx, question in enumerate(WIPEOUT_ROUND_QS):
        wipeout_question_slide = prs.slides.add_slide(slide_layouts.get_by_name("WIPEOUT_ANSWER"))
        wipeout_question_slide.placeholders[10].text = f"Q{idx+1}."
        wipeout_question_slide.placeholders[11].text = question['question']
        wipeout_question_slide.placeholders[12].text = f"A. {question['options'][0]}"
        wipeout_question_slide.placeholders[13].text = f"B. {question['options'][1]}"
        wipeout_question_slide.placeholders[14].text = f"C. {question['options'][2]}"
        wipeout_question_slide.placeholders[15].text = f"D. {question['options'][3]}"
        correct_placeholder_idx = 12
        if question['correct_option'] == "B":
            correct_placeholder_idx = 13
        elif question['correct_option'] == "C":
            correct_placeholder_idx = 14
        elif question['correct_option'] == "D":
            correct_placeholder_idx = 15
        wipeout_question_slide.placeholders[correct_placeholder_idx].fill.solid()
        wipeout_question_slide.placeholders[correct_placeholder_idx].fill.fore_color.rgb = RGBColor(0xE3, 0x09, 0x97)


def create_music_round():
    prs.slides.add_slide(slide_layouts.get_by_name("MUSIC_ROUND_INTRO"))

    def create_song_obj_from_folder(folder):
        try:
            question = next(file_path for file_path in listdir(path.join("questions", "songs", folder)) if re.match("^question\.(mp3|m4a)$", file_path))
            answer = question
            try:
                answer = next(file_path for file_path in listdir(path.join("questions", "songs", folder)) if re.match("^answer\.(mp3|m4a)$", file_path))
            except:
                pass
            img_1 = next(file_path for file_path in listdir(path.join("questions", "songs", folder)) if re.match("^img_1\.(jpe?g|png)$", file_path))
            img_2 = next(file_path for file_path in listdir(path.join("questions", "songs", folder)) if re.match("^img_2\.(jpe?g|png)$", file_path))
            return {"question_track": question, "answer_track": answer, "folder": folder, "img_1": img_1, "img_2": img_2}
        except:
            print(f"Error encountered trying to find right files in folder {folder}")
            exit(1)

    songs = [create_song_obj_from_folder(f) for f in listdir(path.join("questions", "songs"))]
    assert len(songs) == 10, "Need to have 10 folders for audio files in the questions/songs folder"
    for idx, song in enumerate(songs):
        music_question_slide = prs.slides.add_slide(slide_layouts.get_by_name("MUSIC_ROUND_QUESTION"))
        music_question_slide.shapes.add_movie(
            path.join("questions", "songs", song["folder"], song["question_track"]),
            left=Cm(20), top=Cm(-5), width=Cm(2), height=Cm(2),
            mime_type='audio/mp3')
        music_question_slide.placeholders[22].text = f"Track {idx+1}"
    prs.slides.add_slide(slide_layouts.get_by_name("MUSIC_ROUND_ANSWER_BREAK"))
    for idx, song in enumerate(songs):
        music_answer_slide = prs.slides.add_slide(slide_layouts.get_by_name("MUSIC_ROUND_ANSWER"))
        music_answer_slide.shapes.add_movie(
            path.join("questions", "songs", song["folder"], song["answer_track"]),
            left=Cm(20), top=Cm(-5), width=Cm(2), height=Cm(2),
            mime_type='audio/mp3')
        music_answer_slide.placeholders[22].text = f"Track {idx+1}"
        music_answer_slide.placeholders[23].text = song["folder"]
        music_answer_slide.placeholders[24].insert_picture(path.join("questions", "songs", song["folder"], song["img_1"]))
        music_answer_slide.placeholders[25].insert_picture(path.join("questions", "songs", song["folder"], song["img_2"]))


prs.slides.add_slide(slide_layouts.get_by_name("WELCOME"))
prs.slides.add_slide(slide_layouts.get_by_name("RULES"))
create_rounds_preview()
create_worded_round(1)
create_worded_round(2)
create_worded_round(3)
create_picture_round()
create_wipeout_round()
create_music_round()
prs.slides.add_slide(slide_layouts.get_by_name("COUNT_SCORES"))
prs.slides.add_slide(slide_layouts.get_by_name("THANK_YOU"))

prs.save("output.pptx")
